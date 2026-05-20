from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLACK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xFD, 0xF5, 0xE6)
BROWN = RGBColor(0x6B, 0x4F, 0x3A)
DARK_BROWN = RGBColor(0x3E, 0x2C, 0x1C)
ACCENT_ORANGE = RGBColor(0xC8, 0x55, 0x1C)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_BLUE = RGBColor(0x15, 0x65, 0xC0)
LIGHT_BLUE = RGBColor(0xE3, 0xF0, 0xFF)
LIGHT_GREEN = RGBColor(0xDF, 0xF5, 0xDD)
LIGHT_RED = RGBColor(0xFF, 0xE9, 0xE9)
LIGHT_ORANGE = RGBColor(0xFF, 0xF4, 0xE6)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)

def add_bg(slide, color=CREAM):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, corner_radius=None):
    if corner_radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = corner_radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=BLACK, bold_items=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        if bold_items and i in bold_items:
            p.font.bold = True
    return txBox

def add_table_slide(slide, left, top, rows, cols, col_widths, data, header_color=DARK_BROWN, header_text_color=WHITE, font_size=13):
    total_height = Inches(0.4) * rows
    table_shape = slide.shapes.add_table(rows, cols, left, top, sum(col_widths), total_height)
    table = table_shape.table
    for c in range(cols):
        table.columns[c].width = col_widths[c]
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = header_text_color
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = header_color
                else:
                    paragraph.font.color.rgb = BLACK
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHT_GRAY
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
    return table_shape

# ── SLIDE 1: TITLE ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BROWN)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.15), ACCENT_ORANGE)
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.5),
    "Shelby \u00d7 WhatsApp", 54, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
    "Phygital Ordering System", 36, RGBColor(0xE0, 0xC8, 0xA8), False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.8),
    "Same window. Same vibe. Just one less line.", 22, RGBColor(0xBB, 0xA0, 0x80), False, PP_ALIGN.CENTER)
add_shape_bg(slide, Inches(5.5), Inches(5.8), Inches(2.3), Inches(0.05), ACCENT_ORANGE)
add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.6),
    "Owner Presentation  \u2022  May 2026", 16, RGBColor(0x99, 0x88, 0x77), False, PP_ALIGN.CENTER)

# ── SLIDE 2: THE PROBLEM ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "What Problem Are We Solving?", 32, WHITE, True, PP_ALIGN.LEFT)

problems = [
    ("\u26a0\ufe0f  Peak-hour congestion", "The crowd that made you famous is the crowd that's slowing you down. 500+ cups/day through a tiny window."),
    ("\u26a0\ufe0f  The Great Disappointment", "Viral FOMO drives people in, but unannounced 20-minute waits drive negative reviews."),
    ("\u26a0\ufe0f  Repetitive questions drain staff", "\"Are you open?\" \"Vegan milk?\" \"Do you have buns?\" \u2014 heard 100 times a day."),
    ("\u26a0\ufe0f  Menu drift / sold-out confusion", "When Hazelnut runs out, the next 4 customers don't know until they reach the window."),
    ("\u26a0\ufe0f  Rain kills the sidewalk", "Bengaluru dumps rain, the Fourth Space empties, and you have no way to say \"we're still open.\""),
]

for i, (title, desc) in enumerate(problems):
    y = Inches(1.4) + Inches(1.1) * i
    add_shape_bg(slide, Inches(0.6), y, Inches(12), Inches(0.95), WHITE, 0.03)
    add_text_box(slide, Inches(0.9), y + Inches(0.08), Inches(11.5), Inches(0.4), title, 17, ACCENT_RED, True)
    add_text_box(slide, Inches(0.9), y + Inches(0.48), Inches(11.5), Inches(0.45), desc, 14, GRAY)

add_text_box(slide, Inches(0.8), Inches(6.9), Inches(12), Inches(0.5),
    "We are not redesigning your caf\u00e9. We are removing one specific bottleneck: verbal order-taking at the counter during peak hours.",
    14, DARK_BROWN, True, PP_ALIGN.CENTER)

# ── SLIDE 3: THE CORE IDEA ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "The Core Idea \u2014 One QR, One Chat", 32, WHITE, True, PP_ALIGN.LEFT)

steps_data = [
    ("1", "SCAN", "Customer scans QR\nat sidewalk / counter", ACCENT_BLUE, LIGHT_BLUE),
    ("2", "CHAT", "WhatsApp opens.\nBot greets them.", ACCENT_BLUE, LIGHT_BLUE),
    ("3", "ORDER", "Pick items, see ETA,\nget an order code.", ACCENT_ORANGE, LIGHT_ORANGE),
    ("4", "PREP", "Order appears on\nbarista's tablet.", BROWN, LIGHT_ORANGE),
    ("5", "PING", "\"Your order is ready\"\nWhatsApp notification.", ACCENT_GREEN, LIGHT_GREEN),
    ("6", "PICKUP", "Customer steps up\nto window. Done.", ACCENT_GREEN, LIGHT_GREEN),
]
for i, (num, label, desc, border_c, bg_c) in enumerate(steps_data):
    x = Inches(0.5) + Inches(2.1) * i
    y = Inches(1.6)
    box = add_shape_bg(slide, x, y, Inches(1.9), Inches(2.8), bg_c, 0.05)
    box.line.color.rgb = border_c
    box.line.width = Pt(2)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.15), Inches(1.7), Inches(0.5),
        f"{num}. {label}", 18, border_c, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.7), Inches(1.7), Inches(1.8),
        desc, 14, BLACK, False, PP_ALIGN.CENTER)
    if i < 5:
        add_text_box(slide, x + Inches(1.85), y + Inches(1.0), Inches(0.3), Inches(0.5),
            "\u25b6", 20, GRAY, False, PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2), WHITE, 0.03)
add_text_box(slide, Inches(0.8), Inches(4.95), Inches(12), Inches(0.4),
    "That's the whole product.", 20, DARK_BROWN, True)
add_text_box(slide, Inches(0.8), Inches(5.4), Inches(12), Inches(1.4),
    "Everything else we've designed is about making sure this loop survives the messy reality of a real caf\u00e9 \u2014 payment failures, sold-out items, rush hours, rain, internet outages, and customers who just want to talk to a human.",
    15, GRAY)

# ── SLIDE 4: TWO LANES ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "The Most Important Promise: Two Lanes, One Window", 30, WHITE, True, PP_ALIGN.LEFT)

add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.5), LIGHT_BLUE, 0.04)
add_text_box(slide, Inches(1.2), Inches(1.7), Inches(4.8), Inches(0.5),
    "\U0001f4f1  DIGITAL LANE (new)", 22, ACCENT_BLUE, True, PP_ALIGN.CENTER)
lane1_steps = [
    "QR scan at sidewalk / counter / Instagram",
    "WhatsApp chat with Shelby bot",
    "Order placed with honest ETA",
    "Wait anywhere in peace",
    "\U0001f4f2 \"Your order is ready!\"",
]
add_bullet_list(slide, Inches(1.5), Inches(2.3), Inches(4.5), Inches(3.5), lane1_steps, 15, BLACK)

add_shape_bg(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(4.5), LIGHT_ORANGE, 0.04)
add_text_box(slide, Inches(7.4), Inches(1.7), Inches(4.8), Inches(0.5),
    "\U0001f9cd  PHYSICAL LANE (unchanged)", 22, ACCENT_ORANGE, True, PP_ALIGN.CENTER)
lane2_steps = [
    "Walk up to the window",
    "Verbal order to barista",
    "Pay cash or UPI directly",
    "Wait at the window",
    "Pick up when called",
]
add_bullet_list(slide, Inches(7.5), Inches(2.3), Inches(4.5), Inches(3.5), lane2_steps, 15, BLACK)

add_shape_bg(slide, Inches(3.5), Inches(6.2), Inches(6.3), Inches(0.9), ACCENT_GREEN, 0.05)
add_text_box(slide, Inches(3.5), Inches(6.3), Inches(6.3), Inches(0.7),
    "\U0001f3e0  Both lanes converge at the PICKUP WINDOW", 18, WHITE, True, PP_ALIGN.CENTER)

# ── SLIDE 5: NOTHING CHANGES FOR WALK-INS ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "Your Physical Caf\u00e9 Does Not Change", 32, WHITE, True, PP_ALIGN.LEFT)

no_change = [
    ("Cash-only customers?", "Walk up to the window like always."),
    ("No WhatsApp?", "Walk up to the window like always."),
    ("Internet drops?", "Walk up to the window like always."),
    ("Just wants to chat with the barista?", "Walk up to the window like always."),
    ("Prefers the old way?", "Walk up to the window like always."),
]
for i, (q, a) in enumerate(no_change):
    y = Inches(1.5) + Inches(1.0) * i
    add_shape_bg(slide, Inches(1), y, Inches(11.3), Inches(0.85), WHITE, 0.03)
    add_text_box(slide, Inches(1.3), y + Inches(0.1), Inches(4.5), Inches(0.6), q, 18, DARK_BROWN, True)
    add_text_box(slide, Inches(6), y + Inches(0.1), Inches(6), Inches(0.6), a, 18, ACCENT_GREEN, True)

add_text_box(slide, Inches(1), Inches(6.7), Inches(11), Inches(0.5),
    "The WhatsApp lane is additive. It thins the crowd. It never replaces how Shelby works today.",
    16, ACCENT_ORANGE, True, PP_ALIGN.CENTER)

# ── SLIDE 6: HONEST WAIT TIMES ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "Promise #1: Honest Wait Times", 32, WHITE, True, PP_ALIGN.LEFT)

eta_data = [
    ["Kitchen Load", "What the Bot Tells the Customer", "Result"],
    ["< 5 orders in prep", "\"Ready in ~8 minutes\"", "Quick, happy customer"],
    ["5\u201315 orders in prep", "\"Ready in ~10\u201312 minutes\"", "Realistic expectation set"],
    ["> 15 orders (rush)", "\"Rush hour \u2014 ready in 18\u201322 min\"", "No surprise, no disappointment"],
]
add_table_slide(slide, Inches(1), Inches(1.5), 4, 3,
    [Inches(3.5), Inches(5), Inches(3)], eta_data, font_size=15)

add_shape_bg(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(1.5), LIGHT_GREEN, 0.04)
add_text_box(slide, Inches(1.3), Inches(4.4), Inches(10.7), Inches(0.5),
    "\"The Great Disappointment\" is impossible.", 22, ACCENT_GREEN, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.3), Inches(5.0), Inches(10.7), Inches(0.5),
    "We tell people the truth before they pay. They agree to the wait. No surprise 25-minute queue.", 16, BLACK, False, PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(1), Inches(6.0), Inches(11.3), Inches(1.0), WHITE, 0.03)
add_text_box(slide, Inches(1.3), Inches(6.1), Inches(10.7), Inches(0.8),
    "The thresholds (5 / 15) are dials you can tune on the dashboard \u2014 no engineer needed.", 15, GRAY, False, PP_ALIGN.CENTER)

# ── SLIDE 7: WHOLE-ORDER DELIVERY ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "Promise #2: Whole-Order Delivery", 32, WHITE, True, PP_ALIGN.LEFT)

add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0), LIGHT_GREEN, 0.04)
add_text_box(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.5),
    "\u2705  How it works", 22, ACCENT_GREEN, True)
right_steps = [
    "Customer orders: 1 Coffee + 1 Korean Bun",
    "Coffee is ready first (4 min)",
    "Coffee HELD \u2014 bun still cooking",
    "Bun ready (2 more min)",
    "Staff taps READY (whole order assembled)",
    "ONE ping: \"Your order is ready!\"",
    "Customer walks to window ONCE",
]
add_bullet_list(slide, Inches(1.3), Inches(2.3), Inches(4.8), Inches(4.0), right_steps, 15, BLACK)

add_shape_bg(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5.0), LIGHT_RED, 0.04)
add_text_box(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(0.5),
    "\u274c  Why NOT per-item", 22, ACCENT_RED, True)
wrong_steps = [
    "Customer crowds window for coffee",
    "Waits 5 min for bun to arrive",
    "Blocks the next customer",
    "Staff manages half-pickups",
    "Defeats First-Come-First-Serve",
    "More chaos, not less",
]
add_bullet_list(slide, Inches(7.5), Inches(2.3), Inches(4.8), Inches(4.0), wrong_steps, 15, ACCENT_RED)

# ── SLIDE 8: FIVE KILL SWITCHES ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "Promise #3: Five Kill Switches You Control", 30, WHITE, True, PP_ALIGN.LEFT)

switches = [
    ["\U0001f6d1  Item Toggle", "Hazelnut syrup is out", "One tap \u2192 gone from WhatsApp menu instantly"],
    ["\u2614  Rain Mode", "Bengaluru rain starts", "Bot tells everyone: \"we're still open, grab an umbrella!\""],
    ["\u23f1\ufe0f  ETA Factor", "Rush feels heavier than usual", "Push ETAs longer \u2014 manage expectations"],
    ["\u23f8\ufe0f  Pause Digital Lane", "Kitchen is in chaos", "Bot says: \"Walk to the window\" \u2014 total fallback"],
    ["\U0001f91d  Manual Handoff", "Chat got weird / complex", "Take over the conversation yourself"],
]
for i, (name, when, action) in enumerate(switches):
    y = Inches(1.4) + Inches(1.1) * i
    colors = [LIGHT_BLUE, LIGHT_ORANGE, LIGHT_GREEN, LIGHT_RED, LIGHT_BLUE]
    add_shape_bg(slide, Inches(0.6), y, Inches(12), Inches(0.95), colors[i], 0.03)
    add_text_box(slide, Inches(0.9), y + Inches(0.1), Inches(3), Inches(0.7), name, 16, DARK_BROWN, True)
    add_text_box(slide, Inches(4), y + Inches(0.1), Inches(3.5), Inches(0.7), when, 14, GRAY)
    add_text_box(slide, Inches(7.8), y + Inches(0.1), Inches(4.5), Inches(0.7), action, 14, BLACK, True)

add_text_box(slide, Inches(0.8), Inches(6.9), Inches(12), Inches(0.5),
    "No engineer phone call required for any operational lever.", 16, ACCENT_ORANGE, True, PP_ALIGN.CENTER)

# ── SLIDE 9: GRACEFUL DEGRADATION ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "Promise #4: The System Can Break. Your Caf\u00e9 Won't.", 28, WHITE, True, PP_ALIGN.LEFT)

fail_data = [
    ["If This Breaks...", "This Still Works"],
    ["Internet goes down", "Walk-in counter (verbal orders, cash/UPI)"],
    ["WhatsApp provider outage", "Walk-in counter + orders queue and replay when back"],
    ["Razorpay / UPI fails", "\"Pay at counter\" is the default anyway"],
    ["AI / FAQ engine crashes", "Bot falls back to staff handoff"],
    ["Dashboard goes offline", "Staff takes verbal orders; queued orders replay on return"],
    ["Customer taps Confirm 3x", "Idempotency \u2192 ONE order created, never three"],
]
add_table_slide(slide, Inches(1.5), Inches(1.5), 7, 2,
    [Inches(4.5), Inches(6)], fail_data, font_size=15)

add_shape_bg(slide, Inches(3), Inches(5.8), Inches(7), Inches(0.8), ACCENT_GREEN, 0.05)
add_text_box(slide, Inches(3), Inches(5.9), Inches(7), Inches(0.6),
    "\u2615  Coffee never stops being made.", 22, WHITE, True, PP_ALIGN.CENTER)

# ── SLIDE 10: NO AI IN YOUR MONEY ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "Promise #5: No AI In Your Money", 32, WHITE, True, PP_ALIGN.LEFT)

ai_rules = [
    ("AI is used for ONE thing only:", "answering FAQs from your approved knowledge base."),
    ("Every price:", "calculated by deterministic database math. Always."),
    ("Every total:", "calculated by deterministic database math. Always."),
    ("Every availability check:", "happens TWICE \u2014 when added to cart AND at confirmation."),
    ("The system will never:", "invent a \u20b9350 Hazelnut Cold Coffee."),
]
for i, (label, desc) in enumerate(ai_rules):
    y = Inches(1.5) + Inches(0.95) * i
    add_shape_bg(slide, Inches(1), y, Inches(11.3), Inches(0.8), WHITE, 0.03)
    add_text_box(slide, Inches(1.3), y + Inches(0.1), Inches(4), Inches(0.6), label, 16, DARK_BROWN, True)
    add_text_box(slide, Inches(5.3), y + Inches(0.1), Inches(6.8), Inches(0.6), desc, 16, BLACK)

add_shape_bg(slide, Inches(1), Inches(6.3), Inches(11.3), Inches(0.8), LIGHT_RED, 0.04)
add_text_box(slide, Inches(1.3), Inches(6.4), Inches(10.7), Inches(0.6),
    "AI handles conversations. Deterministic code handles money. These two systems never cross.", 16, ACCENT_RED, True, PP_ALIGN.CENTER)

# ── SLIDE 11: THE MENU (from boards) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "Your Menu \u2014 Digitized from the Boards", 30, WHITE, True, PP_ALIGN.LEFT)

menu_left = [
    ["Category", "Items", "Price Range"],
    ["Milk Tea", "Rose, Ginger, Normal, Masala, Vanilla", "\u20b930\u2013\u20b940"],
    ["Black Tea", "Mint Black, Lemon Honey, Clove Special", "\u20b925\u2013\u20b935"],
    ["Milk Coffee", "Shelby Signature, Chocolate, Vanilla, Caramel, Hazelnut", "\u20b950\u2013\u20b970"],
    ["Black Coffee", "Black, Jaggery Blast, Cinnamon-yana", "\u20b925"],
    ["Special", "Horlicks, Boost, Badam Milk, Hot Chocolate, +Rocks", "\u20b935\u2013\u20b9150"],
]
add_table_slide(slide, Inches(0.5), Inches(1.3), 6, 3,
    [Inches(2), Inches(6.5), Inches(2)], menu_left, font_size=12)

menu_right = [
    ["Category", "Items", "Price Range"],
    ["Iced Tea", "Lemon, Peach, Passion Fruit, Elder Flower", "\u20b980\u2013\u20b995"],
    ["Cold Coffee", "Premium, Hazelnut, Vanilla, Caramel, Irish", "\u20b9150\u2013\u20b9180"],
    ["Mojito", "Watermelon, Passion Fruit, Strawberry, Orange, Mango, Virgin, Kala Khatta", "\u20b990"],
    ["Soda / Smoothy", "Lime Soda, Mango Pulpy Smoothy", "\u20b950\u2013\u20b9120"],
    ["Bakery (TBC)", "Korean Cream Cheese Bun \u2014 Classic / Mushroom", "TBD"],
]
add_table_slide(slide, Inches(0.5), Inches(4.2), 6, 3,
    [Inches(2), Inches(6.5), Inches(2)], menu_right, font_size=12)

add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
    "Must-try shortcut on WhatsApp: Coffees \u2022 Hot Chocolate \u2022 Lemon Honey (mirrors your board)", 13, GRAY, False, PP_ALIGN.CENTER)

# ── SLIDE 12: WHAT CHANGES DAY 1 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "Day 1 of Pilot \u2014 What Literally Changes", 30, WHITE, True, PP_ALIGN.LEFT)

change_data = [
    ["Role", "Before Pilot", "During Pilot"],
    ["Customer", "Stand in line, guess wait time", "Optional: scan QR, see ETA, get pinged when ready"],
    ["Order-taker", "Verbal at window, repeats menu 100\u00d7/day", "Verbal only for walk-ins; FAQ deflected to bot"],
    ["Baristas", "Hear orders shouted", "See orders on tablet, batch tea/coffee naturally"],
    ["Owner (you)", "Guess where the bottleneck is", "Dashboard: counter time, ETA delta, channel split"],
    ["Walk-ins / cash", "Normal service", "UNCHANGED \u2014 walk up and order exactly like today"],
]
add_table_slide(slide, Inches(0.5), Inches(1.4), 6, 3,
    [Inches(2.5), Inches(4.5), Inches(5.5)], change_data, font_size=14)

add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.5),
    "Only physical change: add a small QR sticker that says \"Skip the line \u2014 order on WhatsApp\"",
    17, ACCENT_ORANGE, True, PP_ALIGN.CENTER)

# ── SLIDE 13: 4-WEEK ROADMAP ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "4-Week Timeline", 32, WHITE, True, PP_ALIGN.LEFT)

weeks = [
    ("WEEK 1", "Foundation", "Schema, webhooks, event logging,\ndashboard shell + auth", "Nothing on the floor yet", LIGHT_BLUE),
    ("WEEK 2", "Order Engine", "Menu browse, cart, order creation,\nKanban + realtime + toggles", "Internal testing with\nfake orders", LIGHT_ORANGE),
    ("WEEK 3", "AI + Hardening", "FAQ retrieval, handoff, templates,\nidempotency + retries", "Shadow pilot: staff watches\ndashboard during off-peak", LIGHT_GREEN),
    ("WEEK 4", "PILOT", "QR sticker goes up.\nWeekend peak window.", "We sit beside you\nand watch it work.", LIGHT_RED),
]
for i, (week, theme, behind, visible, bg_c) in enumerate(weeks):
    x = Inches(0.5) + Inches(3.15) * i
    y = Inches(1.4)
    add_shape_bg(slide, x, y, Inches(2.95), Inches(5.2), bg_c, 0.04)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.15), Inches(2.75), Inches(0.4),
        week, 20, DARK_BROWN, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.55), Inches(2.75), Inches(0.4),
        theme, 16, ACCENT_ORANGE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(1.1), Inches(2.75), Inches(0.3),
        "Behind the scenes:", 12, GRAY, True)
    add_text_box(slide, x + Inches(0.1), y + Inches(1.4), Inches(2.75), Inches(1.5),
        behind, 13, BLACK)
    add_text_box(slide, x + Inches(0.1), y + Inches(3.2), Inches(2.75), Inches(0.3),
        "What you see:", 12, ACCENT_GREEN, True)
    add_text_box(slide, x + Inches(0.1), y + Inches(3.5), Inches(2.75), Inches(1.2),
        visible, 13, BLACK)

# ── SLIDE 14: SUCCESS METRICS ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "How We Measure Success", 32, WHITE, True, PP_ALIGN.LEFT)

metrics_data = [
    ["Metric", "Target", "Why It Matters to You"],
    ["Counter dwell time", "\u221225%", "Less crowding, fewer encroachment complaints"],
    ["Peak orders via WhatsApp", "\u226520% by week 4", "Real adoption, not a toy"],
    ["WhatsApp order error rate", "<3%", "Trust in the system"],
    ["Bot response time (median)", "<10 seconds", "Customers feel attended to instantly"],
    ["FAQ auto-resolved", "\u226570%", "Staff freed for making coffee, not answering questions"],
    ["ETA promised vs actual", "\u2264\u00b120%", "Brand reputation protected from \"Great Disappointment\""],
]
add_table_slide(slide, Inches(0.8), Inches(1.4), 7, 3,
    [Inches(3.5), Inches(3), Inches(5)], metrics_data, font_size=14)

add_text_box(slide, Inches(0.8), Inches(5.8), Inches(12), Inches(0.5),
    "If any metric misses after 2 weeks of pilot, we adjust before scaling \u2014 not after.",
    16, ACCENT_ORANGE, True, PP_ALIGN.CENTER)

# ── SLIDE 15: PILOT COST ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "Pilot Running Cost", 32, WHITE, True, PP_ALIGN.LEFT)

cost_data = [
    ["Item", "Monthly (approx)"],
    ["WhatsApp BSP (Interakt / AiSensy)", "\u20b91,000 \u2013 \u20b92,500"],
    ["WhatsApp conversation fees (Meta)", "\u20b9500 \u2013 \u20b91,500"],
    ["Supabase (DB + Realtime)", "Free tier covers pilot"],
    ["Hosting (Vercel + Render)", "Free / minimal"],
    ["LLM for FAQ", "\u20b9500 \u2013 \u20b91,000"],
    ["Razorpay", "2% per txn (optional, Phase 3)"],
    ["TOTAL (pilot)", "\u20b92,000 \u2013 \u20b95,000 / month"],
]
add_table_slide(slide, Inches(2), Inches(1.5), 8, 2,
    [Inches(5.5), Inches(3.5)], cost_data, font_size=15)

add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.8),
    "Unit economics improve as Shelby grows (stack scales linearly with volume).\nExcludes development time.", 14, GRAY, False, PP_ALIGN.CENTER)

# ── SLIDE 16: RISKS ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "Risks We Are Carrying (No Surprises)", 30, WHITE, True, PP_ALIGN.LEFT)

risk_data = [
    ["Risk", "Our Plan"],
    ["Kitchen is the real bottleneck, not counter", "We measure both separately. Channel shift relieves counter regardless."],
    ["Customers don't adopt WhatsApp ordering", "Target is only 20% in week 4 \u2014 small, then grow."],
    ["Staff resistance", "Dashboard mirrors the existing assembly line. Train 1 person first."],
    ["AI gives a wrong answer", "Retrieval-only from approved docs. Low confidence = human handoff."],
    ["Fake / no-show orders", "Cart cap \u20b9200 unpaid. Above that = digital pay required."],
    ["Provider lock-in", "Adapter layer \u2192 swap providers in days, not weeks."],
]
add_table_slide(slide, Inches(0.5), Inches(1.4), 7, 2,
    [Inches(5), Inches(7.5)], risk_data, font_size=14)

# ── SLIDE 17: NOT DOING YET ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "What We Are Explicitly NOT Doing (Yet)", 30, WHITE, True, PP_ALIGN.LEFT)

not_doing = [
    ("Delivery (Swiggy/Zomato style)", "Not for MVP"),
    ("Loyalty programs / coupons", "Not for MVP"),
    ("Marketing broadcasts to past customers", "Phase 2 (needs opt-in compliance)"),
    ("Full POS replacement", "Never \u2014 this complements, doesn't replace"),
    ("Voice ordering", "Never \u2014 defeats the calm of WhatsApp"),
    ("Multi-language AI", "Phase 2"),
]
for i, (item, reason) in enumerate(not_doing):
    y = Inches(1.4) + Inches(0.9) * i
    add_shape_bg(slide, Inches(1), y, Inches(11.3), Inches(0.75), LIGHT_RED if "Never" in reason else WHITE, 0.03)
    add_text_box(slide, Inches(1.3), y + Inches(0.1), Inches(5.5), Inches(0.5), item, 16, BLACK, True)
    add_text_box(slide, Inches(7), y + Inches(0.1), Inches(5), Inches(0.5), reason, 16, GRAY)

add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
    "Smaller MVP = clearer pilot = faster learning = less wasted money.", 17, ACCENT_GREEN, True, PP_ALIGN.CENTER)

# ── SLIDE 18: THE DEEPER BET ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "The Deeper Bet: This Is Also a Data Engine", 30, WHITE, True, PP_ALIGN.LEFT)

data_points = [
    "\U0001f4ca  Which items sell out fastest \u2014 and exactly when",
    "\u23f0  Which hours are your real bottlenecks (not guesses \u2014 data)",
    "\U0001f504  Which customers come back (by phone, anonymized)",
    "\u2753  Which questions to put on a permanent sign vs. answer in chat",
    "\U0001f35e  Whether the Korean Bun is worth keeping (or pivoting to the next trend item)",
    "\U0001f4b0  Revenue per peak hour from the digital lane vs. walk-in",
]
for i, point in enumerate(data_points):
    y = Inches(1.5) + Inches(0.85) * i
    add_shape_bg(slide, Inches(1), y, Inches(11.3), Inches(0.7), WHITE, 0.03)
    add_text_box(slide, Inches(1.3), y + Inches(0.08), Inches(10.8), Inches(0.5), point, 16, BLACK)

add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
    "You can't get any of this from a verbal counter today.", 18, ACCENT_ORANGE, True, PP_ALIGN.CENTER)

# ── SLIDE 19: DECISION TABLE ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BROWN)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
    "Decision Table \u2014 Your Sign-Off to Start Week 1", 28, WHITE, True, PP_ALIGN.LEFT)

decision_data = [
    ["Decision", "Default", "Your Call"],
    ["Build with WhatsApp Cloud API via BSP?", "Yes (Interakt / AiSensy)", "\u2610"],
    ["Pilot: Pickup + FAQ + Pay-at-counter only?", "Yes", "\u2610"],
    ["Add Razorpay digital payments in pilot?", "No (add in Phase 3)", "\u2610"],
    ["Add dine-in table ordering in pilot?", "No (add in Phase 1.5)", "\u2610"],
    ["Cart cap for unpaid orders?", "\u20b9200", "\u2610 Tune: ___"],
    ["Rush ETA threshold?", ">15 orders in PREPARING", "\u2610 Tune: ___"],
    ["Korean Buns on bot menu from day 1?", "Yes if available", "\u2610 Confirm price: ___"],
]
add_table_slide(slide, Inches(0.8), Inches(1.3), 8, 3,
    [Inches(5.5), Inches(4), Inches(2.5)], decision_data, font_size=14)

add_shape_bg(slide, Inches(2.5), Inches(5.8), Inches(8), Inches(1.0), ACCENT_GREEN, 0.05)
add_text_box(slide, Inches(2.5), Inches(5.9), Inches(8), Inches(0.8),
    "If you nod on these defaults, Week 1 begins.", 22, WHITE, True, PP_ALIGN.CENTER)

# ── SLIDE 20: CLOSING ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BROWN)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.15), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
    "The One Question to Ask:", 36, RGBColor(0xE0, 0xC8, 0xA8), False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(1.2),
    "\"What's the smallest thing we can do\nin Week 4 to prove this works?\"", 32, WHITE, True, PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(2.0), RGBColor(0x50, 0x3C, 0x28), 0.04)
add_text_box(slide, Inches(1.8), Inches(4.4), Inches(9.7), Inches(1.6),
    "Stick the QR up at 6pm Saturday. Watch the tablet.\nIf even 30 of the 500 cups that night come through WhatsApp,\nwith <3% errors and an honest ETA \u2014 we have a working channel.\n\nThen we scale it across weekday peaks, evenings, dine-in, payments.",
    18, RGBColor(0xE0, 0xC8, 0xA8), False, PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(5.5), Inches(6.6), Inches(2.3), Inches(0.05), ACCENT_ORANGE)
add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
    "Shelby \u00d7 WhatsApp  \u2022  Same window. Same vibe. Just one less line.", 14, RGBColor(0x99, 0x88, 0x77), False, PP_ALIGN.CENTER)

prs.save("deliverables/Shelby_WhatsApp_Owner_Deck.pptx")
print("DONE: deliverables/Shelby_WhatsApp_Owner_Deck.pptx")
print(f"Slides: {len(prs.slides)}")
